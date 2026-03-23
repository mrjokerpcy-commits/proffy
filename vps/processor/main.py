from fastapi import FastAPI, UploadFile, File, HTTPException
from pptx import Presentation
from docx import Document
import io, tempfile, os

app = FastAPI()

@app.get("/health")
def health():
    return {"status": "ok"}

@app.post("/extract/pptx")
async def extract_pptx(file: UploadFile = File(...)):
    data = await file.read()
    try:
        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            title = ""
            content = []
            notes = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if not text:
                        continue
                    if shape.shape_type == 13 or (hasattr(shape, "placeholder_format") and shape.placeholder_format and shape.placeholder_format.idx == 0):
                        title = text
                    else:
                        content.append(text)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            slides.append({"slide_number": i, "title": title, "content": content, "notes": notes})
        return slides
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/extract/docx")
async def extract_docx(file: UploadFile = File(...)):
    data = await file.read()
    try:
        doc = Document(io.BytesIO(data))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return {"text": text}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))
