from fastapi import FastAPI, UploadFile, File, HTTPException, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional, List
import os
import tempfile
import json
from dotenv import load_dotenv

from .pipeline import ProcessingPipeline
from .models import ProcessRequest, ProcessResult

load_dotenv()

app = FastAPI(title="StudyAI Processor", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

pipeline = ProcessingPipeline()


@app.get("/health")
async def health():
    return {"status": "ok"}


@app.post("/process/file", response_model=ProcessResult)
async def process_file(
    file: UploadFile = File(...),
    university: Optional[str] = None,
    department: Optional[str] = None,
    course: Optional[str] = None,
    professor: Optional[str] = None,
    year: Optional[str] = None,
    doc_type: Optional[str] = None,
):
    """Process an uploaded file: extract text, chunk, embed, store in Qdrant."""
    content = await file.read()
    result = await pipeline.process(
        content=content,
        filename=file.filename or "unknown",
        metadata={
            "university": university,
            "department": department,
            "course": course,
            "professor": professor,
            "year": year,
            "type": doc_type,
        },
    )
    return result


@app.post("/process/url")
async def process_url(request: ProcessRequest, background_tasks: BackgroundTasks):
    """Process a file from a URL (Google Drive, web) in the background."""
    background_tasks.add_task(pipeline.process_url, request.url, request.metadata)
    return {"status": "queued", "url": request.url}


class SlideData(BaseModel):
    slide_number: int
    title: str
    content: List[str]
    notes: str


@app.post("/extract/pptx", response_model=List[SlideData])
async def extract_pptx(file: UploadFile = File(...)):
    """Extract slides from a PPTX file. Returns slide-by-slide content with slide numbers.
    Used by the upload route to preserve slide_number for exact citations.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise HTTPException(status_code=501, detail="python-pptx not installed")

    content = await file.read()
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    try:
        prs = Presentation(tmp_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            title = ""
            body_parts = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                # Title placeholder types: 13 (CENTER_TITLE), 15 (TITLE)
                if shape.shape_type in (13, 15) or (hasattr(shape, "placeholder_format") and
                        shape.placeholder_format is not None and
                        shape.placeholder_format.idx == 0):
                    title = text
                else:
                    body_parts.append(text)
            notes = ""
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    notes = notes_frame.text.strip()
            if title or body_parts:
                slides.append(SlideData(
                    slide_number=i + 1,
                    title=title,
                    content=body_parts,
                    notes=notes,
                ))
        return slides
    finally:
        os.unlink(tmp_path)


@app.post("/extract/pdf")
async def extract_pdf_text(file: UploadFile = File(...)):
    """Extract raw text from a PDF using pdfplumber. No page limit. Returns text + page count."""
    try:
        import pdfplumber
    except ImportError:
        raise HTTPException(status_code=501, detail="pdfplumber not installed")

    content = await file.read()
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    try:
        with pdfplumber.open(tmp_path) as pdf:
            pages_text = []
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text and text.strip():
                    pages_text.append(f"[Page {i + 1}]\n{text.strip()}")
            full_text = "\n\n".join(pages_text)
            page_count = len(pdf.pages)
        return {"text": full_text, "pages": page_count}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PDF extraction failed: {str(e)}")
    finally:
        os.unlink(tmp_path)
