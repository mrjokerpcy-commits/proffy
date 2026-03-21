import io
from pptx import Presentation


async def extract_pptx(content: bytes, filename: str) -> str:
    """Extract text from PowerPoint slides."""
    prs = Presentation(io.BytesIO(content))
    text_parts = []

    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            text_parts.append(f"[Slide {i + 1}]\n" + "\n".join(slide_texts))

    return "\n\n".join(text_parts)
